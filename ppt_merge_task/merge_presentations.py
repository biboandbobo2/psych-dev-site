from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import copy
import io
import sys

def copy_slide(source_slide, target_prs):
    # Find layout
    layout_name = source_slide.slide_layout.name
    layout = None
    for l in target_prs.slide_layouts:
        if l.name == layout_name:
            layout = l
            break
    if not layout:
        layout = target_prs.slide_layouts[0] # Fallback
        
    new_slide = target_prs.slides.add_slide(layout)
    
    # Remove existing placeholders to avoid duplication
    for shape in list(new_slide.shapes):
        sp = shape.element
        sp.getparent().remove(sp)
        
    # Copy shapes
    for shape in source_slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            # Handle picture
            try:
                blob = shape.image.blob
                new_slide.shapes.add_picture(io.BytesIO(blob), shape.left, shape.top, shape.width, shape.height)
            except Exception as e:
                print(f"Error copying picture: {e}")
        elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            print(f"Warning: Group shape found. Images inside group might be broken.")
            new_el = copy.deepcopy(shape.element)
            new_slide.shapes._spTree.insert_element_before(new_el, 'p:extLst')
        else:
            # Handle other shapes (Text, AutoShape, etc.)
            new_el = copy.deepcopy(shape.element)
            new_slide.shapes._spTree.insert_element_before(new_el, 'p:extLst')

def merge():
    path1 = "/Users/aleksejzykov/Desktop/архив сайта/35-45.pptx"
    path2 = "/Users/aleksejzykov/Desktop/архив сайта/45-65.pptx"
    output_path = "/Users/aleksejzykov/Desktop/архив сайта/40-65_merged.pptx"
    
    print("Loading presentations...")
    try:
        prs1 = Presentation(path1)
        prs2 = Presentation(path2)
    except Exception as e:
        print(f"Error loading presentations: {e}")
        sys.exit(1)
    
    # 1. Modify Title of prs1 (Slide 1)
    print("Modifying title...")
    if len(prs1.slides) > 0:
        title_slide = prs1.slides[0]
        if title_slide.shapes.title:
            title_slide.shapes.title.text = "Возрастная психология и психология развития\n40-65 лет"
        
    # 2. Remove unwanted slides from prs1
    # "Взрослость (32-42 года)" and "От 34 до 37 лет"
    print("Removing unwanted slides...")
    slides_to_remove = []
    for i, slide in enumerate(prs1.slides):
        text = ""
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text
        if "32-42 года" in text or "От 34 до 37 лет" in text:
            slides_to_remove.append(i)
            
    # Remove in reverse
    for i in sorted(slides_to_remove, reverse=True):
        print(f"Removing slide {i+1}")
        # Delete slide logic
        xml_slides = prs1.slides._sldIdLst
        slides = list(xml_slides)
        xml_slides.remove(slides[i])
        
    temp_path = "/Users/aleksejzykov/Desktop/архив сайта/temp_base.pptx"
    print(f"Saving temp base to {temp_path}...")
    prs1.save(temp_path)
    
    print("Reloading temp base...")
    prs1 = Presentation(temp_path)
    
    # 3. Append slides from prs2 (skipping first)
    print("Appending slides from second presentation...")
    for i in range(1, len(prs2.slides)):
        print(f"Copying slide {i+1} from prs2...")
        copy_slide(prs2.slides[i], prs1)
        
    print(f"Saving to {output_path}...")
    prs1.save(output_path)
    
    # Cleanup temp
    import os
    if os.path.exists(temp_path):
        os.remove(temp_path)
    print("Done.")

if __name__ == "__main__":
    merge()

from pptx import Presentation
import sys
import os

def analyze_ppt(path):
    print(f"--- Analyzing {os.path.basename(path)} ---")
    try:
        prs = Presentation(path)
    except Exception as e:
        print(f"Error opening {path}: {e}")
        return

    for i, slide in enumerate(prs.slides):
        print(f"Slide {i+1}:")
        # Get title
        if slide.shapes.title:
            print(f"  Title: {slide.shapes.title.text}")
        
        # Get all text
        text = []
        has_images = False
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                text.append(shape.text)
            if shape.shape_type == 13: # PICTURE
                has_images = True
        
        if has_images:
            print("  [CONTAINS IMAGES]")
        
        # Join and truncate for display
        full_text = " | ".join(text).replace('\n', ' ')
        print(f"  Content: {full_text[:200]}...")
        print("-" * 20)

if __name__ == "__main__":
    if len(sys.argv) > 1:
        files = sys.argv[1:]
    else:
        files = [
            "/Users/aleksejzykov/Desktop/архив сайта/35-45.pptx",
            "/Users/aleksejzykov/Desktop/архив сайта/45-65.pptx"
        ]
    for f in files:
        analyze_ppt(f)
